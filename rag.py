from flask import Flask, request, jsonify
from langchain_community.document_loaders.pdf import PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate
from langchain.chains import create_retrieval_chain
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain.prompts import ChatPromptTemplate, PromptTemplate
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_mongodb import MongoDBAtlasVectorSearch
from pymongo import MongoClient
from langchain_huggingface import HuggingFaceEndpoint
from dotenv import load_dotenv
from flask_cors import CORS, cross_origin
from pptx import Presentation
from pptx.util import Inches
import google.generativeai as genai
import os

load_dotenv()

embeddings = HuggingFaceEmbeddings(model_name="mixedbread-ai/mxbai-embed-large-v1", encode_kwargs={'precision': 'binary'})

# Vector Store
one_bit_vectorstore = FAISS.load_local("OS", embeddings, allow_dangerous_deserialization=True)
retriever = one_bit_vectorstore.as_retriever(search_kwargs={"k": 10})  # Increased context retrieval

one_bit_vectorstore_ml = FAISS.load_local("ML", embeddings, allow_dangerous_deserialization=True)
retriever_ml = one_bit_vectorstore_ml.as_retriever(search_kwargs={"k": 10})  # Increased context retrieval

app = Flask(__name__)
cors = CORS(app)
app.config['CORS_HEADERS'] = 'Content-Type'

# Configure the Gemini API
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=GOOGLE_API_KEY)

# def create_slide(presentation, title, content_points):
#     slide_layout = presentation.slide_layouts[1]  # Use the layout with title and content
#     slide = presentation.slides.add_slide(slide_layout)
    
#     title_placeholder = slide.shapes.title
#     content_placeholder = slide.placeholders[1]
    
#     title_placeholder.text = title
    
#     for point in content_points:
#         p = content_placeholder.text_frame.add_paragraph()
#         p.text = point
#         p.level = 0  # Bullet point level

# def generate_presentation(slide_data, template_path, output_path):
#     presentation = Presentation(template_path)
    
#     for slide_info in slide_data:
#         title = slide_info['title']
#         content_points = slide_info['content']
        
#         while content_points:
#             slide = presentation.slides.add_slide(presentation.slide_layouts[1])
#             title_placeholder = slide.shapes.title
#             content_placeholder = slide.placeholders[1]
            
#             title_placeholder.text = title
            
#             # Add points to the slide until it is full
#             remaining_points = []
#             for point in content_points:
#                 p = content_placeholder.text_frame.add_paragraph()
#                 p.text = point
#                 p.level = 0  # Bullet point level
                
#                 # Check if the content exceeds the slide's capacity
#                 if content_placeholder.text_frame.fit_text():
#                     remaining_points.append(point)
#                     content_placeholder.text_frame.text = content_placeholder.text_frame.text.rsplit('\n', 1)[0]
#                     break
            
#             content_points = remaining_points
    
#     presentation.save(output_path)

def is_comparative_or_exploratory_question(question):
    """
    Detect if the question is comparative or exploratory
    """
    comparative_keywords = [
        'difference', 'compare', 'contrast', 'best', 'better', 'vs', 'versus', 
        'what is', 'explain', 'describe', 'how', 'why', 'characteristics', 
        'features', 'advantages', 'disadvantages'
    ]
    
    # Convert question to lowercase for case-insensitive matching
    lower_question = question.lower()
    
    # Check if any comparative keywords are in the question
    return any(keyword in lower_question for keyword in comparative_keywords)

@app.route("/os", methods=['POST'])
def hello_world():
    content = request.json

    # Gemini model configuration
    model = genai.GenerativeModel('gemini-1.5-pro')

    # Improved template with flexible answering strategy
    template = """
    Context: {context}
    
    Question: {question}
    
Guidelines for Response:

1. Responses will be strictly based on the information available in the vector database. If the context provides direct information relevant to the question, the response will use that information verbatim, citing the exact document name and section as the source.

2. For comparative or exploratory questions, responses will be structured with a comprehensive explanation based only on relevant documents in the vector database. Each key point will be attributed to specific sections of the referenced documents.

3. If no relevant information is available in the vector database, the response will clearly state that the information is not present. General knowledge or external sources will not be used.

4. Answers will be precise, structured, and directly relevant to the query. Citations will be provided in a reference format, ensuring clarity and credibility.

5. Every claim or piece of information in the response will be accompanied by an inline reference number [1], [2], etc., which corresponds to the citation list at the end. The "Sources" section will list the exact document name and section from the vector database used to generate the response.

6. If no information is found in the vector database, the response will explicitly state:"No relevant information is available in the vector database."

7. Every response will end with a structured "Sources:" section that lists the exact document names and sections from which the information was retrieved.
    """

    # Add basic validation for question content
    if 'question' not in content or not content['question'].strip():
        return {"answer": "Please provide a valid question."}, 400

    # Add error handling for the API call
    try:
        # Retrieve context first
        context_docs = retriever.invoke(content['question'])
        context_text = "\n".join([doc.page_content for doc in context_docs])

        # Determine if it's a comparative or exploratory question
        is_comparative = is_comparative_or_exploratory_question(content['question'])

        # Prepare the full prompt
        full_prompt = template.format(context=context_text, question=content['question'])

        # If it's a comparative or exploratory question, add more context
        if is_comparative and context_text:
            full_prompt += "\n\nNote: This is a comparative or exploratory question. Provide a comprehensive analysis using available context and broader knowledge."

        # Generate response using Gemini
        response = model.generate_content(full_prompt)
        
        return {"answer": response.text}
    except Exception as e:
        return {"answer": f"An error occurred while processing your request. Please try again with a more specific question.", "error": str(e)}, 500

@app.route("/ml", methods=['POST'])
def ml():
    content = request.json

    # Gemini model configuration
    model = genai.GenerativeModel('gemini-1.5-pro')

    # Improved template with flexible answering strategy
    template = """
    Context: {context}
    
    Question: {question}
    
Guidelines for Response:

1. Responses will be strictly based on the information available in the vector database. If the context provides direct information relevant to the question, the response will use that information verbatim, citing the exact document name and section as the source.

2. For comparative or exploratory questions, responses will be structured with a comprehensive explanation based only on relevant documents in the vector database. Each key point will be attributed to specific sections of the referenced documents.

3. If no relevant information is available in the vector database, the response will clearly state that the information is not present. General knowledge or external sources will not be used.

4. Answers will be precise, structured, and directly relevant to the query. Citations will be provided in a reference format, ensuring clarity and credibility.

5. Every claim or piece of information in the response will be accompanied by an inline reference number [1], [2], etc., which corresponds to the citation list at the end. The "Sources" section will list the exact document name and section from the vector database used to generate the response.

6. If no information is found in the vector database, the response will explicitly state:"No relevant information is available in the vector database."

7. Every response will end with a structured "Sources:" section that lists the exact document names and sections from which the information was retrieved.
    """

    # Add basic validation for question content
    if 'question' not in content or not content['question'].strip():
        return {"answer": "Please provide a valid question."}, 400

    # Add error handling for the API call
    try:
        # Retrieve context first
        context_docs = retriever_ml.invoke(content['question'])
        context_text = "\n".join([doc.page_content for doc in context_docs])

        # Determine if it's a comparative or exploratory question
        is_comparative = is_comparative_or_exploratory_question(content['question'])

        # Prepare the full prompt
        full_prompt = template.format(context=context_text, question=content['question'])

        # If it's a comparative or exploratory question, add more context
        if is_comparative and context_text:
            full_prompt += "\n\nNote: This is a comparative or exploratory question. Provide a comprehensive analysis using available context and broader knowledge."

        # Generate response using Gemini
        response = model.generate_content(full_prompt)
        
        return {"answer": response.text}
    except Exception as e:
        return {"answer": f"An error occurred while processing your request. Please try again with a more specific question.", "error": str(e)}, 500


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=os.getenv("PORT"))